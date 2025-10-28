from io import BytesIO

from utils import texttools

from .file_processors import is_docx, is_pptx


# File processor functions following Single Responsibility Principle
def process_docx(file_content: BytesIO) -> str:
    """Process DOCX files and return extracted text."""
    import docx

    doc = docx.Document(file_content)
    text = "\n\n".join(p.text for p in doc.paragraphs if p.text)
    return texttools.normalize_text(text)


def process_pptx(file_content: BytesIO) -> str:
    """Process PPTX files and return extracted text."""
    import pptx

    prs = pptx.Presentation(file_content)
    texts: list[str] = []
    for slide in prs.slides:
        texts.extend(
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
    return texttools.normalize_text("\n\n".join(texts))


def process_direct_file(file_content: BytesIO, file_type: str) -> str:
    """Process direct file and return extracted text."""
    if is_docx(file_type):
        return process_docx(file_content)
    elif is_pptx(file_type):
        return process_pptx(file_content)
    else:
        return ""
